# slide_translator/pptx_processor.py
from pptx.presentation import Presentation
from pptx.shapes.group import GroupShape
from typing import Dict, Iterator, Any, Tuple

def _iter_shapes(shapes) -> Iterator[Any]:
    """
    Recursively iterate over shapes, including those within groups.
    """
    for shape in shapes:
        if isinstance(shape, GroupShape):
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape

def extract_paragraphs_with_run_mapping(prs: Presentation) -> Tuple[Dict[str, str], Dict[str, list]]:
    """
    Groups text runs by paragraph for better translation context, while maintaining 
    individual run mapping for precise formatting preservation.

    Args:
        prs: A python-pptx Presentation object.

    Returns:
        A tuple containing:
        - A dictionary mapping {paragraph_id: combined_paragraph_text}.
        - A dictionary mapping {paragraph_id: list_of_run_objects_in_paragraph}.
    """
    paragraph_texts: Dict[str, str] = {}
    paragraph_runs: Dict[str, list] = {}

    for slide_idx, slide in enumerate(prs.slides):
        # Process shapes on the slide
        for shape_idx, shape in enumerate(_iter_shapes(slide.shapes)):
            if shape.has_text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    para_id = f"s{slide_idx}_sh{shape_idx}_p{para_idx}"
                    
                    # Combine all runs in this paragraph into a single text
                    combined_text = ""
                    runs_in_paragraph = []
                    
                    for run in paragraph.runs:
                        combined_text += run.text
                        runs_in_paragraph.append(run)
                    
                    if combined_text.strip():  # Only add non-empty paragraphs
                        paragraph_texts[para_id] = combined_text
                        paragraph_runs[para_id] = runs_in_paragraph
            
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            para_id = f"s{slide_idx}_sh{shape_idx}_t_r{row_idx}_c{cell_idx}_p{para_idx}"
                            
                            combined_text = ""
                            runs_in_paragraph = []
                            
                            for run in paragraph.runs:
                                combined_text += run.text
                                runs_in_paragraph.append(run)
                            
                            if combined_text.strip():
                                paragraph_texts[para_id] = combined_text
                                paragraph_runs[para_id] = runs_in_paragraph

        # Process notes slide
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for para_idx, paragraph in enumerate(slide.notes_slide.notes_text_frame.paragraphs):
                para_id = f"s{slide_idx}_notes_p{para_idx}"
                
                combined_text = ""
                runs_in_paragraph = []
                
                for run in paragraph.runs:
                    combined_text += run.text
                    runs_in_paragraph.append(run)
                
                if combined_text.strip():
                    paragraph_texts[para_id] = combined_text
                    paragraph_runs[para_id] = runs_in_paragraph

    return paragraph_texts, paragraph_runs
