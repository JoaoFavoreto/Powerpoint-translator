from typing import Dict, List, Tuple, Optional
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame, _Paragraph, _Run
import tempfile
import os
import logging

from core.models import TextRun
from core.exceptions import FileProcessingError

logger = logging.getLogger(__name__)

class PPTXService:
    """Serviço para processamento de arquivos PowerPoint"""
    
    def __init__(self):
        self.run_mapping: Dict[str, _Run] = {}
        self.text_mapping: Dict[str, str] = {}
        self.paragraph_mapping: Dict[str, any] = {}  # Novo: para parágrafos
        self.presentation: Optional[Presentation] = None
    
    def extract_texts(self, file_path: str) -> Dict[str, str]:
        """Extrai textos do PowerPoint mantendo referências dos runs"""
        
        try:
            logger.debug(f"🔍 Iniciando extração de textos de: {file_path}")
            
            # Carregar e manter referência da apresentação
            self.presentation = Presentation(file_path)
            self.run_mapping = {}
            self.text_mapping = {}
            self.paragraph_mapping = {}  # Novo: mapear parágrafos
            
            run_counter = 0
            paragraph_counter = 0
            
            logger.debug(f"📄 Slides encontrados: {len(self.presentation.slides)}")
            
            for slide_idx, slide in enumerate(self.presentation.slides):
                logger.debug(f"📑 Processando slide {slide_idx + 1}")
                
                # Processar shapes principais (incluindo títulos, caixas de texto, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        run_counter, paragraph_counter = self._extract_from_text_frame(
                            shape.text_frame, slide_idx, run_counter, paragraph_counter)
                    
                    # Processar shapes de grupo (podem conter textos)
                    elif hasattr(shape, "shapes"):
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text_frame") and sub_shape.text_frame:
                                run_counter, paragraph_counter = self._extract_from_text_frame(
                                    sub_shape.text_frame, slide_idx, run_counter, paragraph_counter)
                    
                    # Processar tabelas
                    elif shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                if hasattr(cell, "text_frame") and cell.text_frame:
                                    run_counter, paragraph_counter = self._extract_from_text_frame(
                                        cell.text_frame, slide_idx, run_counter, paragraph_counter,
                                        f"table_{row_idx}_{col_idx}")
                    
                    # Processar charts (gráficos) se tiverem texto
                    elif hasattr(shape, "has_chart") and shape.has_chart:
                        try:
                            chart = shape.chart
                            if hasattr(chart, "chart_title") and chart.chart_title:
                                if hasattr(chart.chart_title, "text_frame") and chart.chart_title.text_frame:
                                    run_counter, paragraph_counter = self._extract_from_text_frame(
                                        chart.chart_title.text_frame, slide_idx, run_counter, paragraph_counter, "chart_title")
                        except Exception as e:
                            logger.debug(f"Erro ao processar gráfico: {e}")
                
                # Processar notas do slide
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                        run_counter, paragraph_counter = self._extract_from_text_frame(
                            notes_slide.notes_text_frame, slide_idx, run_counter, paragraph_counter, "notes")
            
            logger.debug(f"✅ Extração concluída:")
            logger.debug(f"   - Runs extraídos: {len(self.run_mapping)}")
            logger.debug(f"   - Parágrafos extraídos: {len(self.paragraph_mapping)}")
            logger.debug(f"   - Total de textos: {len(self.text_mapping)}")
            
            if self.text_mapping:
                logger.debug("📝 Primeiros textos encontrados:")
                for i, (key, text) in enumerate(list(self.text_mapping.items())[:3]):
                    logger.debug(f"   {key}: '{text[:50]}...'")
            
            return self.text_mapping
            
        except Exception as e:
            logger.error(f"❌ Erro ao extrair textos: {e}")
            raise FileProcessingError(f"Erro ao extrair textos do PowerPoint: {e}")
    
    def _extract_from_text_frame(self, text_frame, slide_idx, run_counter, paragraph_counter, prefix=""):
        """Extrai textos de um text_frame específico"""
        
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            # Verificar se o parágrafo tem texto
            paragraph_text = paragraph.text.strip()
            if not paragraph_text:
                continue
            
            # Para parágrafos com bullet points, múltiplos runs, ou que sejam títulos/headings
            is_heading = len(paragraph_text) <= 50 and len(paragraph.runs) == 1  # Provável título
            has_bullets = paragraph.level > 0  # Bullet point
            has_multiple_runs = len(paragraph.runs) > 1  # Múltiplos runs
            
            if has_bullets or has_multiple_runs or is_heading:
                paragraph_id = f"para_{paragraph_counter}"
                self.paragraph_mapping[paragraph_id] = paragraph
                self.text_mapping[paragraph_id] = paragraph_text
                paragraph_counter += 1
                logger.debug(f"Parágrafo detectado: {paragraph_id} ({'título' if is_heading else 'bullet' if has_bullets else 'múltiplo'})")
            else:
                # Para runs únicos simples, extrair individualmente  
                for run in paragraph.runs:
                    if run.text.strip():
                        run_id = f"run_{run_counter}"
                        self.run_mapping[run_id] = run
                        self.text_mapping[run_id] = run.text
                        run_counter += 1
        
        return run_counter, paragraph_counter
    
    def get_text_runs_info(self) -> List[TextRun]:
        """Retorna informações detalhadas sobre os runs de texto"""
        
        runs_info = []
        
        for run_id, run in self.run_mapping.items():
            text_run = TextRun(
                id=run_id,
                text=run.text,
                is_bold=run.font.bold if run.font.bold is not None else False,
                is_italic=run.font.italic if run.font.italic is not None else False,
                font_name=run.font.name,
                font_size=run.font.size.pt if run.font.size else None
            )
            runs_info.append(text_run)
        
        return runs_info
    
    def apply_translations(self, translations: Dict[str, str]) -> None:
        """Aplica as traduções ao documento PowerPoint preservando formatação."""
        if not self.presentation:
            raise FileProcessingError("Apresentação não carregada")
        
        try:
            logger.debug(f"Aplicando {len(translations)} traduções")
            applied_count = 0
            
            # Aplicar traduções por parágrafo (preserva formatação de bullet points)
            for para_id, paragraph in self.paragraph_mapping.items():
                if para_id in translations:
                    # Preservar formatação: limpar runs e recriar com formatação do primeiro run
                    if paragraph.runs:
                        # Salvar formatação do primeiro run
                        first_run = paragraph.runs[0]
                        font_name = first_run.font.name
                        font_size = first_run.font.size
                        font_bold = first_run.font.bold
                        font_italic = first_run.font.italic
                        font_color = first_run.font.color.rgb if first_run.font.color.type == 1 else None
                        
                        # Limpar texto atual
                        paragraph.clear()
                        
                        # Adicionar novo run com formatação preservada
                        new_run = paragraph.add_run()
                        new_run.text = translations[para_id]
                        if font_name:
                            new_run.font.name = font_name
                        if font_size:
                            new_run.font.size = font_size
                        if font_bold is not None:
                            new_run.font.bold = font_bold
                        if font_italic is not None:
                            new_run.font.italic = font_italic
                        if font_color:
                            new_run.font.color.rgb = font_color
                    else:
                        # Fallback se não houver runs
                        paragraph.text = translations[para_id]
                    
                    applied_count += 1
                    logger.debug(f"Parágrafo traduzido: {para_id}")
            
            # Aplicar traduções por run (preserva formatação original)
            for run_id, run in self.run_mapping.items():
                if run_id in translations:
                    # Verificar se este run não faz parte de um parágrafo já traduzido
                    run_in_paragraph = any(
                        run in paragraph.runs 
                        for paragraph in self.paragraph_mapping.values()
                    )
                    if not run_in_paragraph:
                        # Preservar formatação original do run
                        run.text = translations[run_id]
                        applied_count += 1
                        logger.debug(f"Run traduzido: {run_id}")
            
            logger.info(f"{applied_count} traduções aplicadas ao PowerPoint")
            self.translated = True
        except Exception as e:
            logger.error(f"Erro ao aplicar traduções: {e}")
            raise FileProcessingError(f"Erro ao aplicar traduções: {e}")
    
    def save_presentation(self, original_path: str, output_path: str) -> str:
        """Salva a apresentação traduzida"""
        
        try:
            if self.presentation is None:
                raise FileProcessingError("Nenhuma apresentação carregada. Execute extract_texts primeiro.")
            
            # Salvar a apresentação com as modificações já aplicadas
            self.presentation.save(output_path)
            return output_path
            
        except Exception as e:
            raise FileProcessingError(f"Erro ao salvar apresentação: {e}")
    
    def create_temp_file(self, original_filename: str) -> str:
        """Cria um arquivo temporário para a apresentação traduzida"""
        
        try:
            # Extrair nome base e extensão
            name, ext = os.path.splitext(original_filename)
            
            # Criar arquivo temporário
            temp_file = tempfile.NamedTemporaryFile(
                delete=False,
                suffix=f"_translated{ext}",
                prefix=f"{name}_"
            )
            
            return temp_file.name
            
        except Exception as e:
            raise FileProcessingError(f"Erro ao criar arquivo temporário: {e}")
    
    def get_presentation_stats(self, file_path: str) -> Dict[str, int]:
        """Retorna estatísticas da apresentação"""
        
        try:
            prs = Presentation(file_path)
            
            total_slides = len(prs.slides)
            total_shapes = 0
            total_text_runs = 0
            total_characters = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    total_shapes += 1
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    total_text_runs += 1
                                    total_characters += len(run.text)
            
            return {
                "slides": total_slides,
                "shapes": total_shapes,
                "text_runs": total_text_runs,
                "characters": total_characters
            }
            
        except Exception as e:
            raise FileProcessingError(f"Erro ao obter estatísticas: {e}")
